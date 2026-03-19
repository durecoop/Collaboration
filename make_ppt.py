from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 색상
PURPLE = RGBColor(102, 126, 234)
DARK = RGBColor(51, 51, 51)
GRAY = RGBColor(153, 153, 153)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(220, 38, 38)
BLUE = RGBColor(79, 70, 229)
GREEN = RGBColor(34, 197, 94)
ORANGE = RGBColor(217, 119, 6)
PINK = RGBColor(240, 147, 251)
SKY = RGBColor(79, 172, 254)
LIGHT_BG = RGBColor(249, 250, 251)
RED_BG = RGBColor(254, 226, 226)
BLUE_BG = RGBColor(224, 231, 255)
GREEN_BG = RGBColor(220, 252, 231)
YELLOW_BG = RGBColor(254, 243, 199)
PURPLE_BG = RGBColor(237, 233, 254)

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def set_text(shape, text, size=18, color=DARK, bold=False, align=PP_ALIGN.CENTER):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf

def text(slide, left, top, w, h, txt, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf

def multiline(slide, left, top, w, h, lines):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (t, s, c, b) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = t
        p.font.size = Pt(s)
        p.font.color.rgb = c
        p.font.bold = b
        p.space_after = Pt(6)
    return tf

def bg_white(slide):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = WHITE

def bg_purple(slide):
    fill = slide.background.fill; fill.solid(); fill.fore_color.rgb = PURPLE

def arrow_right(slide, left, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, Inches(0.8), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PURPLE
    shape.line.fill.background()
    return shape

def arrow_down(slide, left, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, Inches(0.5), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PURPLE
    shape.line.fill.background()
    return shape


# ── 슬라이드 1: 표지 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_purple(slide)
text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
    "🤝 같이 개발하는 법", 44, WHITE, True, PP_ALIGN.CENTER)
text(slide, Inches(1), Inches(3.2), Inches(11), Inches(1),
    "Commit · Push · Pull 쉽게 이해하기", 26, RGBColor(200, 210, 255), False, PP_ALIGN.CENTER)

for i, (name, color) in enumerate([("지원", RGBColor(90,111,214)), ("소영", PINK), ("슬기", SKY)]):
    c = add_circle(slide, Inches(4.5) + Inches(i*1.8), Inches(4.5), Inches(1.2), color)
    set_text(c, name[0], 32, WHITE, True)

text(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.6),
    "초등학생도 이해할 수 있는 Git 협업 가이드 📖", 20, RGBColor(200,210,255), False, PP_ALIGN.CENTER)


# ── 슬라이드 2: 비유 - 공유 노트 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_white(slide)
text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
    "🤔 같이 개발한다는 게 뭐야?", 36, DARK, True)

box = add_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5), LIGHT_BG)
text(slide, Inches(1.2), Inches(2), Inches(4.8), Inches(0.7),
    "학교에서 비유하면...", 22, DARK, True)
multiline(slide, Inches(1.2), Inches(2.8), Inches(4.8), Inches(3.5), [
    ("📓  반 공용 노트가 있어요", 20, DARK, False),
    ("", 10, DARK, False),
    ("✏️  각자 자기 파트를 써요", 20, DARK, False),
    ("", 10, DARK, False),
    ("📤  다 쓰면 노트에 끼워 넣어요", 20, DARK, False),
    ("", 10, DARK, False),
    ("📥  다른 친구가 쓴 것도 볼 수 있어요", 20, DARK, False),
])

box2 = add_rect(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(5), BLUE_BG)
text(slide, Inches(7.2), Inches(2), Inches(4.8), Inches(0.7),
    "개발에서는...", 22, BLUE, True)
multiline(slide, Inches(7.2), Inches(2.8), Inches(4.8), Inches(3.5), [
    ("📓  GitHub = 공용 노트", 20, DARK, False),
    ("", 10, DARK, False),
    ("✏️  각자 컴퓨터에서 코드 작성", 20, DARK, False),
    ("", 10, DARK, False),
    ("📤  Push = 노트에 끼워 넣기", 20, DARK, False),
    ("", 10, DARK, False),
    ("📥  Pull = 친구가 쓴 거 가져오기", 20, DARK, False),
])

text(slide, Inches(6.05), Inches(3.8), Inches(1), Inches(0.8),
    "→", 40, PURPLE, True, PP_ALIGN.CENTER)


# ── 슬라이드 3: 전체 흐름 한 그림 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_white(slide)
text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
    "📋 전체 흐름을 한 그림으로 보면", 36, DARK, True)

# 내 컴퓨터
my_pc = add_rect(slide, Inches(0.5), Inches(1.5), Inches(3.5), Inches(5.2), LIGHT_BG)
text(slide, Inches(0.5), Inches(1.6), Inches(3.5), Inches(0.6),
    "💻 내 컴퓨터", 22, DARK, True, PP_ALIGN.CENTER)

step1 = add_rect(slide, Inches(0.8), Inches(2.4), Inches(2.9), Inches(1), GREEN_BG)
set_text(step1, "① 코드 작성\n파일을 수정해요", 16, DARK, False)
arrow_down(slide, Inches(1.9), Inches(3.5))
step2 = add_rect(slide, Inches(0.8), Inches(4.2), Inches(2.9), Inches(1), YELLOW_BG)
set_text(step2, "② Commit\n\"저장\" 버튼 누르기", 16, DARK, False)
arrow_down(slide, Inches(1.9), Inches(5.3))
step3 = add_rect(slide, Inches(0.8), Inches(6), Inches(2.9), Inches(0.7), BLUE_BG)
set_text(step3, "③ Push → 올리기", 16, BLUE, True)

# 화살표 Push
arrow_right(slide, Inches(4.2), Inches(6.1))

# GitHub (가운데)
gh_box = add_rect(slide, Inches(5.2), Inches(1.5), Inches(3), Inches(5.2), PURPLE_BG)
text(slide, Inches(5.2), Inches(1.6), Inches(3), Inches(0.6),
    "☁️ GitHub", 22, PURPLE, True, PP_ALIGN.CENTER)
text(slide, Inches(5.2), Inches(2.3), Inches(3), Inches(0.5),
    "(인터넷 저장소)", 16, GRAY, False, PP_ALIGN.CENTER)

gh_inner = add_rect(slide, Inches(5.5), Inches(3), Inches(2.4), Inches(3), WHITE)
multiline(slide, Inches(5.7), Inches(3.2), Inches(2), Inches(2.5), [
    ("📁 CLAUDE.md", 14, DARK, False),
    ("📁 README.md", 14, DARK, False),
    ("📁 docs/", 14, DARK, False),
    ("📁 src/", 14, DARK, False),
    ("", 8, DARK, False),
    ("모든 파일이", 13, GRAY, False),
    ("여기에 보관돼요", 13, GRAY, False),
])

# 화살표 Pull
shape = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(8.4), Inches(6.1), Inches(0.8), Inches(0.5))
shape.fill.solid(); shape.fill.fore_color.rgb = GREEN; shape.line.fill.background()

# 팀원 컴퓨터
team_pc = add_rect(slide, Inches(9.4), Inches(1.5), Inches(3.5), Inches(5.2), LIGHT_BG)
text(slide, Inches(9.4), Inches(1.6), Inches(3.5), Inches(0.6),
    "💻 팀원 컴퓨터", 22, DARK, True, PP_ALIGN.CENTER)

t1 = add_rect(slide, Inches(9.7), Inches(2.4), Inches(2.9), Inches(0.7), GREEN_BG)
set_text(t1, "④ Pull ← 받기", 16, GREEN, True)
arrow_down(slide, Inches(10.8), Inches(3.2))
t2 = add_rect(slide, Inches(9.7), Inches(3.9), Inches(2.9), Inches(1), GREEN_BG)
set_text(t2, "⑤ 코드 작성\n자기 파트 수정", 16, DARK, False)
arrow_down(slide, Inches(10.8), Inches(5))
t3 = add_rect(slide, Inches(9.7), Inches(5.7), Inches(2.9), Inches(1), YELLOW_BG)
set_text(t3, "⑥ Commit + Push\n저장하고 올리기", 16, DARK, False)


# ── 슬라이드 4: Commit ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_white(slide)
text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
    "💾 Commit = 내 컴퓨터에 저장", 36, DARK, True)

box = add_rect(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(2), YELLOW_BG)
multiline(slide, Inches(1.5), Inches(2), Inches(10), Inches(1.7), [
    ("게임에서 \"세이브\" 하는 거랑 똑같아요!", 24, DARK, True),
    ("", 10, DARK, False),
    ("게임 중간중간 세이브하면 죽어도 그 지점부터 다시 하잖아요?", 20, GRAY, False),
    ("코드도 중간중간 Commit(저장)하면 잘못되어도 돌아갈 수 있어요!", 20, GRAY, False),
])

# 예시
text(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(0.6),
    "Commit 할 때 메모를 남겨요 (뭘 했는지 적기)", 20, DARK, True)

commits = [
    ("📝  \"로그인 화면 만들었음\"", "지원 · 오늘 오전 10시", BLUE_BG),
    ("📝  \"할 일 목록 디자인 수정\"", "소영 · 오늘 오후 2시", RGBColor(252, 231, 243)),
    ("📝  \"마감일 기능 추가\"", "슬기 · 오늘 오후 4시", RGBColor(219, 234, 254)),
]
for i, (msg, who, color) in enumerate(commits):
    y = Inches(4.9) + Inches(i * 0.75)
    card = add_rect(slide, Inches(1.5), y, Inches(10), Inches(0.65), color)
    text(slide, Inches(1.8), y + Inches(0.05), Inches(5), Inches(0.55), msg, 18, DARK, True)
    text(slide, Inches(7), y + Inches(0.05), Inches(4), Inches(0.55), who, 15, GRAY, False, PP_ALIGN.RIGHT)

text(slide, Inches(0.8), Inches(7), Inches(11), Inches(0.4),
    "⚠️ Commit은 내 컴퓨터에만 저장! 아직 팀원은 못 봐요", 18, RED, True, PP_ALIGN.CENTER)


# ── 슬라이드 5: Push ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_white(slide)
text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
    "📤 Push = GitHub에 올리기", 36, DARK, True)

box = add_rect(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(2), BLUE_BG)
multiline(slide, Inches(1.5), Inches(2), Inches(10), Inches(1.7), [
    ("카톡 단체방에 사진 올리는 거랑 비슷해요!", 24, DARK, True),
    ("", 10, DARK, False),
    ("내가 찍은 사진(=코드)을 단체방(=GitHub)에 올리면", 20, GRAY, False),
    ("다른 친구들도 볼 수 있게 되는 거예요!", 20, GRAY, False),
])

# 그림: 내 컴퓨터 → GitHub
pc = add_rect(slide, Inches(1.5), Inches(4.3), Inches(3.5), Inches(2.5), LIGHT_BG)
text(slide, Inches(1.5), Inches(4.4), Inches(3.5), Inches(0.6),
    "💻 내 컴퓨터", 20, DARK, True, PP_ALIGN.CENTER)
multiline(slide, Inches(1.8), Inches(5.1), Inches(3), Inches(1.5), [
    ("commit 한 내용들:", 14, GRAY, False),
    ("✅ 로그인 화면 완성", 16, DARK, False),
    ("✅ 버튼 색상 변경", 16, DARK, False),
])

# 큰 화살표
shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.3), Inches(5), Inches(1.8), Inches(1))
shape.fill.solid(); shape.fill.fore_color.rgb = PURPLE; shape.line.fill.background()
set_text(shape, "Push!", 22, WHITE, True)

gh = add_rect(slide, Inches(7.5), Inches(4.3), Inches(4.5), Inches(2.5), PURPLE_BG)
text(slide, Inches(7.5), Inches(4.4), Inches(4.5), Inches(0.6),
    "☁️ GitHub (인터넷)", 20, PURPLE, True, PP_ALIGN.CENTER)
multiline(slide, Inches(7.8), Inches(5.1), Inches(4), Inches(1.5), [
    ("올라온 내용:", 14, GRAY, False),
    ("✅ 로그인 화면 완성", 16, DARK, False),
    ("✅ 버튼 색상 변경", 16, DARK, False),
    ("→ 팀원 모두 볼 수 있음!", 15, PURPLE, True),
])


# ── 슬라이드 6: Pull ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_white(slide)
text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
    "📥 Pull = 팀원 작업 받아오기", 36, DARK, True)

box = add_rect(slide, Inches(0.8), Inches(1.8), Inches(11.5), Inches(2), GREEN_BG)
multiline(slide, Inches(1.5), Inches(2), Inches(10), Inches(1.7), [
    ("카톡 단체방에 새 사진 와서 다운 받는 거랑 같아요!", 24, DARK, True),
    ("", 10, DARK, False),
    ("친구가 올린 사진(=코드)을 내 폰(=컴퓨터)에 받으면", 20, GRAY, False),
    ("최신 상태로 업데이트되는 거예요!", 20, GRAY, False),
])

# 그림: GitHub → 내 컴퓨터
gh = add_rect(slide, Inches(1.5), Inches(4.3), Inches(4), Inches(2.5), PURPLE_BG)
text(slide, Inches(1.5), Inches(4.4), Inches(4), Inches(0.6),
    "☁️ GitHub", 20, PURPLE, True, PP_ALIGN.CENTER)
multiline(slide, Inches(1.8), Inches(5.1), Inches(3.5), Inches(1.5), [
    ("새로 올라온 것:", 14, GRAY, False),
    ("🆕 소영: 디자인 수정", 16, DARK, False),
    ("🆕 슬기: 마감일 기능", 16, DARK, False),
])

shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.8), Inches(5), Inches(1.8), Inches(1))
shape.fill.solid(); shape.fill.fore_color.rgb = GREEN; shape.line.fill.background()
set_text(shape, "Pull!", 22, WHITE, True)

pc = add_rect(slide, Inches(7.9), Inches(4.3), Inches(4.5), Inches(2.5), LIGHT_BG)
text(slide, Inches(7.9), Inches(4.4), Inches(4.5), Inches(0.6),
    "💻 내 컴퓨터", 20, DARK, True, PP_ALIGN.CENTER)
multiline(slide, Inches(8.2), Inches(5.1), Inches(4), Inches(1.5), [
    ("받아온 결과:", 14, GRAY, False),
    ("✅ 내 코드 + 소영 코드 + 슬기 코드", 16, DARK, False),
    ("→ 전부 합쳐진 최신 상태!", 15, GREEN, True),
])


# ── 슬라이드 7: 실제 하는 순서 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_white(slide)
text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
    "⏰ 매일 이렇게 하면 돼요!", 36, DARK, True)

steps = [
    ("1", "Pull 받기", "시작할 때 팀원이 한 거 받아오기", "📥", GREEN_BG, GREEN),
    ("2", "코드 작성", "내가 맡은 부분 작업하기", "✏️", LIGHT_BG, DARK),
    ("3", "Commit", "내 컴퓨터에 저장하기", "💾", YELLOW_BG, ORANGE),
    ("4", "Push", "GitHub에 올려서 팀원과 공유", "📤", BLUE_BG, BLUE),
]

for i, (num, title, desc, icon, bg_color, title_color) in enumerate(steps):
    x = Inches(0.5) + Inches(i * 3.2)
    card = add_rect(slide, x, Inches(1.8), Inches(2.8), Inches(3.5), bg_color)

    num_circle = add_circle(slide, x + Inches(1), Inches(2), Inches(0.8), title_color)
    set_text(num_circle, num, 28, WHITE, True)

    text(slide, x, Inches(3), Inches(2.8), Inches(0.6),
        f"{icon} {title}", 24, title_color, True, PP_ALIGN.CENTER)
    text(slide, x, Inches(3.7), Inches(2.8), Inches(1),
        desc, 16, GRAY, False, PP_ALIGN.CENTER)

    if i < 3:
        arrow_right(slide, x + Inches(2.9), Inches(3.3))

# Claude 팁
tip = add_rect(slide, Inches(1), Inches(5.7), Inches(11), Inches(1.3), PURPLE_BG)
multiline(slide, Inches(1.5), Inches(5.9), Inches(10), Inches(1), [
    ('💡 Claude에게 말로 시키면 돼요!', 22, PURPLE, True),
    ('   "pull 해줘" → "커밋하고 푸시해줘" 이 두 마디면 끝!', 20, DARK, False),
])


# ── 슬라이드 8: 충돌이 나면? ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_white(slide)
text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
    "💥 충돌! 같은 부분을 고치면?", 36, DARK, True)

# 상황 설명
box = add_rect(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.5), RED_BG)
multiline(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2), [
    ("지원이가 1번줄을 \"사과\"로 바꾸고, 소영이도 1번줄을 \"바나나\"로 바꿨어요", 20, DARK, False),
    ("둘 다 Push 하면... 컴퓨터가 헷갈려요! 뭘 쓸지 모르니까요 😵", 20, RED, True),
])

# 해결
text(slide, Inches(0.8), Inches(3.5), Inches(11), Inches(0.6),
    "해결 방법:", 24, DARK, True)

solutions = [
    ("1. 예방하기", "서로 다른 파일을 작업해요\n\"나는 로그인, 너는 디자인\"", GREEN_BG),
    ("2. 브랜치 쓰기", "각자 복사본에서 작업하고\n나중에 합쳐요 (다음 장에서!)", BLUE_BG),
    ("3. Claude한테 맡기기", "\"충돌 해결해줘\" 하면\nClaude가 알아서 정리!", PURPLE_BG),
]
for i, (title, desc, color) in enumerate(solutions):
    x = Inches(0.5) + Inches(i * 4.2)
    card = add_rect(slide, x, Inches(4.3), Inches(3.8), Inches(2.5), color)
    text(slide, x + Inches(0.3), Inches(4.5), Inches(3.2), Inches(0.6), title, 20, DARK, True, PP_ALIGN.CENTER)
    text(slide, x + Inches(0.3), Inches(5.2), Inches(3.2), Inches(1.2), desc, 16, GRAY, False, PP_ALIGN.CENTER)


# ── 슬라이드 9: 브랜치 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_white(slide)
text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
    "🌿 브랜치 = 각자 연습장에서 작업", 36, DARK, True)

box = add_rect(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.3), YELLOW_BG)
multiline(slide, Inches(1.5), Inches(1.7), Inches(10), Inches(1.1), [
    ("시험지에 바로 쓰면 지우기 어렵잖아요?", 22, DARK, True),
    ("연습장(브랜치)에서 먼저 쓰고, 잘 됐으면 시험지(main)에 옮기는 거예요!", 20, GRAY, False),
])

# main 줄
main_line = add_rect(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.6), DARK)
set_text(main_line, "main (완성본)", 18, WHITE, True)

# 브랜치 3개
branches = [
    ("지원: 로그인", RGBColor(90,111,214), Inches(3.3)),
    ("소영: 디자인", PINK, Inches(4.5)),
    ("슬기: 마감일기능", SKY, Inches(5.7)),
]
for name, color, top in branches:
    line = add_rect(slide, Inches(3), top, Inches(5), Inches(0.5), color)
    set_text(line, name, 16, WHITE, True)
    # 분기 표시
    small = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.8), top + Inches(0.1), Inches(0.3), Inches(0.3))
    small.fill.solid(); small.fill.fore_color.rgb = color; small.line.fill.background()
    # 합치기 화살표
    arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.2), top + Inches(0.05), Inches(0.6), Inches(0.4))
    arr.fill.solid(); arr.fill.fore_color.rgb = color; arr.line.fill.background()

# 합치기 설명
merge_box = add_rect(slide, Inches(9), Inches(3.8), Inches(3), Inches(2.8), GREEN_BG)
multiline(slide, Inches(9.3), Inches(4), Inches(2.5), Inches(2.5), [
    ("✅ 다 되면", 18, GREEN, True),
    ("main에 합쳐요!", 18, GREEN, True),
    ("", 10, DARK, False),
    ("이걸 Merge", 16, DARK, False),
    ("라고 해요", 16, DARK, False),
])

text(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.5),
    '💡 Claude에게: "브랜치 만들어줘" → 작업 → "main에 머지해줘"', 18, PURPLE, True, PP_ALIGN.CENTER)


# ── 슬라이드 10: 실제 사용법 (Claude) ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_white(slide)
text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
    "🗣️ Claude에게 이렇게 말하세요", 36, DARK, True)

text(slide, Inches(0.8), Inches(1.4), Inches(11), Inches(0.6),
    "어려운 명령어 외울 필요 없어요. 그냥 말하면 돼요!", 22, GRAY, False)

commands = [
    ("시작할 때", '"pull 해줘"', "최신 코드 받아옴", GREEN_BG),
    ("작업 끝나면", '"커밋하고 푸시해줘"', "저장하고 올림", BLUE_BG),
    ("뭐 바뀌었는지 보려면", '"최근 변경사항 보여줘"', "팀원 작업 확인", YELLOW_BG),
    ("내 브랜치 만들려면", '"브랜치 만들어줘"', "연습장 생성", PURPLE_BG),
    ("완성해서 합치려면", '"main에 머지해줘"', "완성본에 합침", GREEN_BG),
    ("충돌 났을 때", '"충돌 해결해줘"', "알아서 정리", RED_BG),
]

for i, (when, say, result, color) in enumerate(commands):
    y = Inches(2.2) + Inches(i * 0.82)
    row = add_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.7), color)
    text(slide, Inches(1.1), y + Inches(0.1), Inches(3), Inches(0.5), when, 17, DARK, True)
    text(slide, Inches(4.2), y + Inches(0.1), Inches(4.5), Inches(0.5), say, 18, DARK, False)
    text(slide, Inches(8.8), y + Inches(0.1), Inches(3.2), Inches(0.5), f"→ {result}", 16, GRAY, False)


# ── 슬라이드 11: 정리 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_white(slide)
text(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(1),
    "📝 3줄 정리!", 36, DARK, True)

summary = [
    ("💾 Commit", "내 컴퓨터에 저장하기", "게임 세이브 같은 거", YELLOW_BG),
    ("📤 Push", "GitHub에 올리기", "카톡에 사진 올리는 거", BLUE_BG),
    ("📥 Pull", "팀원 작업 받아오기", "카톡에서 사진 다운 받는 거", GREEN_BG),
]

for i, (title, desc, analogy, color) in enumerate(summary):
    y = Inches(1.5) + Inches(i * 1.8)
    card = add_rect(slide, Inches(1), y, Inches(11), Inches(1.5), color)
    text(slide, Inches(1.5), y + Inches(0.15), Inches(3), Inches(0.6), title, 28, DARK, True)
    text(slide, Inches(1.5), y + Inches(0.7), Inches(4), Inches(0.6), desc, 20, DARK, False)
    text(slide, Inches(7), y + Inches(0.4), Inches(4.5), Inches(0.6),
        f'"{analogy}"', 20, GRAY, False, PP_ALIGN.CENTER)


# ── 슬라이드 12: 마무리 ──
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_purple(slide)
text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
    "이제 같이 만들어 봐요! 🚀", 44, WHITE, True, PP_ALIGN.CENTER)

text(slide, Inches(1), Inches(3.5), Inches(11), Inches(1.2),
    "Pull → 작업 → Commit → Push\n이것만 기억하면 돼요!", 28, RGBColor(200,210,255), False, PP_ALIGN.CENTER)

for i, (name, color) in enumerate([("지원", RGBColor(90,111,214)), ("소영", PINK), ("슬기", SKY)]):
    c = add_circle(slide, Inches(4.5) + Inches(i*1.8), Inches(5.2), Inches(1.2), color)
    set_text(c, name[0], 32, WHITE, True)


prs.save('D:/2_Projects/automation/Collaboration/docs/Git-협업-가이드.pptx')
print("완료!")
